# app_dashboard/views.py
import os
from django.contrib import messages
from django.contrib.auth.decorators import login_required
from django.shortcuts import get_object_or_404, redirect, render

from actividades.forms import ActividadForm, InformeForm,InformeContenidoForm,PerfilUsuarioForm, TipoActividadForm,EvidenciaForm, UserManagementForm, UserPerfilForm
from actividades.models import Actividad, Informe,Evidencia, InformeFirmado, TipoActividad
from autenticacion.models import PerfilUsuario
from django.contrib.auth.models import Group
from django.http import HttpResponseForbidden, JsonResponse
from django.contrib.auth.models import User
from django.db.models import Q

from django.http import HttpResponse
from django.views import View
from docx import Document
from pptx import Presentation
from .models import AuditoriaInforme, Informe, Actividad

from reportlab.lib.pagesizes import letter
from reportlab.lib.units import inch
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Image, Spacer, PageBreak, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet
from datetime import datetime
from reportlab.lib.enums import TA_CENTER
from reportlab.lib.styles import ParagraphStyle
from reportlab.platypus import Spacer
import pandas as pd
from django.db.models import Count, Q
from docx.shared import Inches,Pt
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from datetime import date, timedelta
from django.template.loader import get_template
from xhtml2pdf import pisa
from django.db.models.functions import TruncMonth
from calendar import month_name
import locale
from django.core.paginator import Paginator
from django.templatetags.static import static
from django.utils.timezone import now
from io import BytesIO
from PIL import Image as PILImage
from django.core.files.base import ContentFile
from docx.oxml.ns import nsdecls
from docx.oxml import parse_xml
from bs4 import BeautifulSoup
from docx.shared import Cm
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor 
from django.views.decorators.cache import never_cache


def get_user_permissions(user):
    permissions = {}
    # Obteniendo todos los permisos del usuario
    for perm in user.get_all_permissions():
        perm_name = perm.split('.')[-1]  # Extraer el nombre del permiso
        permissions[f'can_{perm_name}'] = True
    return permissions

@never_cache
@login_required
def dashboard(request):
    user = request.user  # Usuario autenticado
    permissions = get_user_permissions(user) 
    
    perfil_usuario = PerfilUsuario.objects.get(user=request.user)
    programa_usuario = perfil_usuario.programa
    
    # Establecer el idioma en español
    locale.setlocale(locale.LC_TIME, 'es_ES.UTF-8')  # Para sistemas Linux/Mac
    locale.setlocale(locale.LC_TIME, 'es_ES')       # Para Windows
    
    # Determinar el grupo al que pertenece el usuario
    es_director = user.groups.filter(name='director').exists()
    es_profesor = user.groups.filter(name='profesor').exists()

    # Filtrar actividades según el tipo de usuario
    if es_director:
        actividades = Actividad.objects.filter(creado_por__perfilusuario__programa=user.perfilusuario.programa).distinct()
    elif es_profesor:
        actividades = Actividad.objects.filter(Q(creado_por=user) | Q(colaborador=user)).distinct()
    

    # Datos básicos de actividades
    total_actividades = actividades.count()
    actividades_finalizadas = actividades.filter(estado='Finalizada').count()
    actividades_pendientes = actividades.filter(estado='Pendiente').count()
    actividades_proceso = actividades.filter(estado='En Proceso').count()

    # Actividades por estado
    actividades_finalizadas_lista = actividades.filter(estado='Finalizada')
    actividades_pendientes_lista = actividades.filter(estado='Pendiente')
    actividades_proceso_lista = actividades.filter(estado='En Proceso')

    # Actividades con alertas
    hoy = date.today()
    actividades_atrasadas = actividades.filter(Q(estado='Pendiente') | Q(estado='En proceso'), fecha_fin__lt=hoy)
    actividades_proximas_vencer = actividades.filter(
        estado='Pendiente', fecha_fin__range=[hoy, hoy + timedelta(days=7)]
    )

    num_actividades_atrasadas = actividades_atrasadas.count()
    num_actividades_proximas_vencer = actividades_proximas_vencer.count()

    # Informes firmados y no firmados
    informes_firmados = InformeFirmado.objects.filter(estado_firma=True).count()
    informes_no_firmados = InformeFirmado.objects.filter(estado_firma=False).count()

    # Aportes de colaboradores
    colaborador_aportes = User.objects.filter(
        actividades_asignadas__in=actividades
    ).annotate(
        total_informes=Count('informe', distinct=True),
        informes_firmados=Count(
            'informefirmado',
            filter=Q(
                informefirmado__estado_firma=True,
                informefirmado__informe__actividad__in=actividades,
            ),
            distinct=True,
        ),
    )

    # Actividades por tipo
    actividades_por_tipo = actividades.values('tipo__id', 'tipo__nombre').annotate(
        total=Count('id')
    )
    # Filtrar los tipos de actividad que pertenecen al programa del usuario autenticado
    total_tipos_actividades = TipoActividad.objects.filter(programa=programa_usuario)


    # Crear estructura con conteos de actividades por tipo
    total_tipos_actividades_con_conteo = [
        {
            "id": tipo.id,
            "nombre": tipo.nombre,
            "cantidad": actividades_por_tipo.filter(tipo__id=tipo.id).aggregate(total=Count('id'))['total'] or 0,
        }
        for tipo in total_tipos_actividades
    ]
    
    # Serializar datos para gráficos
    tipo_labels = [tipo["nombre"] for tipo in total_tipos_actividades_con_conteo]
    actividad_counts = [tipo["cantidad"] for tipo in total_tipos_actividades_con_conteo]
    
    # Inicializar variables específicas
    colaboradores = []
    

    # Filtrar datos según el tipo de usuario
    if es_director:
        # Mostrar colaboradores del programa del director
        colaboradores = User.objects.filter(
            perfilusuario__programa=user.perfilusuario.programa
        ).annotate(
            total_actividades=Count('actividades_asignadas'),
        )

    informes_pendientes = []
    etiquetas_meses = []
    datos_meses = []
    if es_profesor:
        # 1. Obtener informes pendientes (estado_firma=False)
        informes_pendientes = InformeFirmado.objects.filter(
            estado_firma=False,
            usuario=user
        ).select_related('informe', 'informe__actividad')

        if not informes_pendientes.exists():
            # 2. Agrupar informes firmados (estado_firma=True) por mes
            informes_firmados_por_mes = (
                InformeFirmado.objects.filter(estado_firma=True, usuario=user)
                .annotate(mes=TruncMonth('fecha_firma'))  # Agrupa por mes
                .values('mes')  # Sólo tomamos el campo `mes`
                .annotate(total=Count('id'))  # Contamos informes firmados
                .order_by('mes')
            )

            # 3. Inicializar un diccionario con 0 informes para cada mes
            meses_totales = {i: 0 for i in range(1, 13)}  # Meses de enero (1) a diciembre (12)

            for informe in informes_firmados_por_mes:
                mes = informe['mes'].month  # Obtener el número del mes
                meses_totales[mes] = informe['total']  # Guardar el total de informes en el mes

            # 4. Generar etiquetas y datos para el gráfico
            etiquetas_meses = [month_name[i] for i in range(1, 13)]  # Nombres de los meses
            datos_meses = [meses_totales[i] for i in range(1, 13)] 



    # Contexto
    context = {
        'total_actividades': total_actividades,
        'actividades_finalizadas': actividades_finalizadas,
        'actividades_pendientes': actividades_pendientes,
        'actividades_proceso': actividades_proceso,
        'informes_firmados': informes_firmados,
        'informes_no_firmados': informes_no_firmados,
        'actividades_atrasadas': actividades_atrasadas,
        'actividades_proximas_vencer': actividades_proximas_vencer,
        'num_actividades_atrasadas': num_actividades_atrasadas,
        'num_actividades_proximas_vencer': num_actividades_proximas_vencer,
        'colaborador_aportes': colaborador_aportes,
        'tipo_labels': tipo_labels,
        'actividad_counts': actividad_counts,
        'total_actividades_lista': actividades,
        'actividades_finalizadas_lista': actividades_finalizadas_lista,
        'actividades_proceso_lista': actividades_proceso_lista,
        'actividades_pendientes_lista': actividades_pendientes_lista,
        'total_tipos_actividades': total_tipos_actividades.count(),
        'total_tipos_actividades_lista': total_tipos_actividades_con_conteo,
        'es_director': es_director,
        'es_profesor': es_profesor,
        'colaboradores': colaboradores,
        'informes_pendientes': informes_pendientes,
        'etiquetas_meses': etiquetas_meses,
        'datos_meses': datos_meses,
        **permissions,
    }

    return render(request, 'actividades/dashboard.html', context)


def obtener_arbol_tipo_actividad(request, tipo_id):
    
    try:
        tipo = TipoActividad.objects.get(id=tipo_id)
        actividades = tipo.actividad_set.all()  # Relación entre TipoActividad y Actividades
        
        # Verificar si el usuario pertenece al grupo 'Director'
        es_director = request.user.groups.filter(name='Director').exists()
        
        # Si el usuario es director, mostrar todas las actividades del tipo
        if es_director:
            actividades = tipo.actividad_set.all()
        else:
            # Si no es director, mostrar solo las actividades creadas o asignadas al usuario
            actividades = tipo.actividad_set.filter(
                Q(creado_por=request.user) | Q(colaborador=request.user)
            ).distinct()
        
        # Creación del nodo padre (TipoActividad)
        data = {
            "name": tipo.nombre,
            "children": []  # Inicializamos la lista de nodos hijos
        }

        # Recorremos las actividades relacionadas con el TipoActividad
        # Agregar actividades e informes
        for actividad in actividades:
            actividad_data = {
                "name": actividad.titulo,
                "children": [
                    {"name": informe.titulo, "value": informe.id}
                    for informe in actividad.informes.all()
                ]
            }
            
            # Agregamos la actividad con sus informes al nodo padre (TipoActividad)
            data["children"].append(actividad_data)

        print("Datos del árbol:", data)  # Mostrar los datos en la consola del servidor
        return JsonResponse(data, safe=False)

    except TipoActividad.DoesNotExist:
        return JsonResponse({"error": "Tipo de actividad no encontrado"}, status=404)


    

#############################################                                ACTIVIDADES                      #################################################################
@never_cache
@login_required
def listar_actividad(request):
    user = request.user
    permissions = get_user_permissions(user)
    programa_usuario = user.perfilusuario.programa  # El programa al que pertenece el usuario
    # Obtener solo los tipos de actividad relacionados con el programa del usuario
    tipos_de_actividad = TipoActividad.objects.filter(programa=programa_usuario)
    
    # Verificar si es una solicitud AJAX
    if request.headers.get('x-requested-with') == 'XMLHttpRequest':
        query = request.GET.get('q', '').strip()  # Consulta de búsqueda

        # Filtrar actividades según los permisos del usuario (director o no director)
        if user.groups.filter(name="Director").exists():
            actividades = Actividad.objects.filter(creado_por__perfilusuario__programa=user.perfilusuario.programa).distinct()
        else:
            actividades = Actividad.objects.filter(Q(creado_por=user) | Q(colaborador=user)).distinct()

        # Filtrar según el query de búsqueda (por título o descripción)
        if query:
            actividades = actividades.filter(titulo__icontains=query)
            

        # Serializar los datos de las actividades
        actividades_data = [{
            'id': actividad.id,
            'titulo': actividad.titulo,
            'descripcion': actividad.descripcion,
            'lugar': actividad.lugar,
            'fecha_inicio': actividad.fecha_inicio.strftime('%Y-%m-%d'),
            'fecha_fin': actividad.fecha_fin.strftime('%Y-%m-%d'),
            'estado': actividad.estado,
            'tipo': actividad.tipo.id if actividad.tipo else None,
            'colaborador': [
                {
                    'id':colaborador.id,
                    'nombre':colaborador.first_name,
                    'apellido': colaborador.last_name,
                    'img': colaborador.perfilusuario.img_usuario.url if colaborador.perfilusuario and colaborador.perfilusuario.img_usuario else static('actividades/img/perfil.png'),

                } for colaborador in actividad.colaborador.all()
            ],
            'colaboradores_ids': list(actividad.colaborador.values_list('id', flat=True)) or [],
        } for actividad in actividades]

        # Devolver los datos filtrados como respuesta JSON
        return JsonResponse({'actividades': actividades_data})
    
    # Obtener los filtros desde el formulario
    mes = request.GET.get('mes')
    dia = request.GET.get('dia')
    estado = request.GET.get('estado', '').lower()
    
    # Verificar si el usuario pertenece al grupo "Director"
    if user.groups.filter(name="Director").exists():
        actividades = Actividad.objects.filter(creado_por__perfilusuario__programa=user.perfilusuario.programa).distinct()
    else:
        # Si no es director, solo puede ver actividades creadas por él o en las que es colaborador
        actividades = Actividad.objects.filter(Q(creado_por=user) | Q(colaborador=user)).distinct()

    
    # Filtrar por mes si se ha seleccionado un mes
    if mes:
        actividades = actividades.filter(fecha_inicio__month=mes)

    # Filtrar por día si se ha seleccionado un día
    if dia:
        actividades = actividades.filter(fecha_inicio__day=dia)

    if estado in ['pendiente', 'en proceso', 'finalizada']:
        actividades = actividades.filter(estado=estado)

    # Para cada actividad, determinar si el usuario es el creador o colaborador
    for actividad in actividades:
        actividad.es_creador = actividad.creado_por == user
        actividad.es_colaborador = actividad.colaborador.filter(id=user.id).exists()
        actividad_colaboradores_ids = actividad.colaborador.values_list('id', flat=True)
        actividad.colaboradores_ids = list(actividad_colaboradores_ids)
        actividad.verificar_informes()
        
    # Paginación
    paginator = Paginator(actividades, 20)
    page_number = request.GET.get('page')
    page_obj = paginator.get_page(page_number)

    dias = list(range(1, 32))
    form_actividad = ActividadForm(programa_usuario=programa_usuario)

    context = {
        'actividades': page_obj.object_list,
        'form_actividad': form_actividad,
        'dias': dias,
        'page_obj': page_obj,
        'tipos_de_actividad': tipos_de_actividad,
        **permissions
    }
    return render(request, 'actividades/listar_actividad.html', context)



@never_cache
@login_required
def agregar_actividad(request):
    user = request.user  # Obtener el usuario autenticado
    # Obtener permisos de manera dinámica
    permissions = get_user_permissions(user)
    
    if request.method == 'POST':
        form = ActividadForm(request.POST, request.FILES)  # Asegúrate de incluir request.FILES para las imágenes
        if form.is_valid():
            # Guardar la actividad y asignar el creador
            actividad = form.save(commit=False)
            actividad.creado_por = user
            actividad.save()
            
            # Guardar los colaboradores seleccionados
            colaboradores_ids = request.POST.getlist('colaborador')  # Obtener los IDs de colaboradores seleccionados
            actividad.colaborador.set(colaboradores_ids)  # Establecer la relación muchos a muchos

            # Devolver una respuesta JSON indicando éxito
            return JsonResponse({'mensaje': 'Actividad guardada correctamente.'}, status=200)
        else:
            # En caso de error, devolver errores del formulario
            return JsonResponse({'errores': form.errors}, status=400)
    else:
        print(form.errors)
        
    form = ActividadForm()
    
    contex = {
        'form': form, 
        **permissions
    }

    return render(request, 'actividades/listar_actividad.html', contex)
@never_cache
@login_required
def editar_actividad(request):
    user = request.user  # Obtener el usuario autenticado
    # Obtener permisos de manera dinámica
    permissions = get_user_permissions(user)
    # Asegúrate de que estás usando la clase User para filtrar
    colaboradores = User.objects.filter(groups__name='Profesor')
    
    if request.method == 'POST':
        actividad_id = request.POST.get('actividad_id')
        actividad = get_object_or_404(Actividad, id=actividad_id)
        
        # Creamos el formulario pasando el objeto existente
        form_actividad = ActividadForm(request.POST, request.FILES, instance=actividad)
        
        if form_actividad.is_valid():
            # Guardar el formulario para actualizar la actividad existente
            form_actividad.save()
            
            # Guardar los colaboradores seleccionados
            colaboradores_ids = request.POST.getlist('colaborador')  # Obtener los IDs de colaboradores seleccionados
            actividad.colaborador.set(colaboradores_ids)  # Establecer la relación muchos a muchos
            
            # Retornar una respuesta exitosa (o redirigir como lo prefieras)
            return JsonResponse({'success': True})
        else:
            return redirect(listar_actividad)
    else:
        # Si el método es GET, pre-cargar el formulario con los datos de la actividad a editar
        actividad_id = request.GET.get('actividad_id')
        actividad = get_object_or_404(Actividad, id=actividad_id)
        form_actividad = ActividadForm(instance=actividad)

    context = {
        'form_actividad': form_actividad,
        'actividad': actividad,
        'colaboradores': colaboradores,
        **permissions
    }
    
    return render(request, 'actividades/listar_actividad.html', context)

@never_cache
@login_required
def eliminar_actividad(request, actividad_id):
    actividad = get_object_or_404(Actividad, id=actividad_id)
    if request.method == "POST":
        AuditoriaInforme.objects.filter(informe__actividad=actividad).delete()
        actividad.delete()  # Elimina la actividad
        messages.success(request, "Actividad eliminada con éxito.")
        return redirect("listar_actividad")  # Redirige a la lista de actividades
    return render(request, "actividades/listar_actividad.html", {"actividad": actividad})



@never_cache#############################################                                INFORMES                           #################################################################
@login_required
def listar_informe(request, actividad_id):
    user = request.user
    permissions = get_user_permissions(user)
    
    actividad = get_object_or_404(Actividad, id=actividad_id)
    informes = Informe.objects.filter(actividad=actividad)
    
    # Obtener todos los usuarios asignados a la actividad
    usuarios_actividad = actividad.colaborador.all() 
    
    # Verificar y agregar nuevos registros en InformeFirmado
    for informe in informes:
        for usuario in usuarios_actividad:
            # Crear un registro en InformeFirmado si no existe para este usuario e informe
            informe_firmado, created = InformeFirmado.objects.get_or_create(
                informe=informe,
                usuario=usuario,
                defaults={'estado_firma': False}  # Por defecto, no firmado
            )
    
    # Verificar si hay algún colaborador que pertenezca al grupo 'Director'
    grupo_directores = Group.objects.get(name='Director')  # Asegúrate de que el nombre sea correcto
    directores = usuarios_actividad.filter(groups=grupo_directores)
    
    es_director = grupo_directores in user.groups.all()  # Verifica si el usuario es director
    es_colaborador = user in usuarios_actividad
    
    # Si hay directores, firmar automáticamente todos los informes
    for director in directores:
        for informe in informes:
            informe_firmado, created = InformeFirmado.objects.get_or_create(informe=informe, usuario=director)
            if not informe_firmado.estado_firma:
                informe_firmado.estado_firma = True
                informe_firmado.save()
    
    for informe in informes:
        # Se verifica si el usuario actual ha firmado el informe
        informe.estado_firma_usuario = informe.firmas.filter(usuario=user, estado_firma=True).exists()
        
        # Obtener usuarios que han firmado
        informe.usuarios_firmados = User.objects.filter(id__in=informe.firmas.filter(estado_firma=True).values_list('usuario_id', flat=True)).distinct()

        
        # Obtener todos los usuarios que NO han firmado (estado_firma=False)
        informe.usuarios_no_firmados = User.objects.filter(id__in=informe.firmas.filter(estado_firma=False).values_list('usuario_id', flat=True)).distinct()
        
        # Verificar si el usuario es el creador del informe
        informe.es_creador = informe.usuario == user
        
    actividad.verificar_informes()
        
    form_informe = InformeForm()
    
    context = {
        'actividad': actividad,
        'informes': informes,
        'form_informe': form_informe,
        'es_director': es_director,
        'es_colaborador': es_colaborador,
        **permissions
    }
    
    return render(request, 'actividades/listar_informe.html', context)

@never_cache
@login_required
def firmar_todos_informes(request, actividad_id):
    actividad = get_object_or_404(Actividad, id=actividad_id)
    
    user = request.user
    
    # Verificar que el usuario es colaborador de la actividad
    if user not in actividad.colaborador.all():
        return HttpResponseForbidden("No tienes permiso para firmar estos informes.")
    
    # Obtener todos los informes de la actividad
    informes = actividad.informes.all()
    
    # Firmar todos los informes no firmados por el usuario (profesor)
    for informe in informes:
        informe_firmado, created = InformeFirmado.objects.get_or_create(informe=informe,usuario=user,)
        if not informe_firmado.estado_firma:
            informe_firmado.estado_firma = True
            informe_firmado.save()

    # Verificar si todos los informes están firmados y actualizar el estado de la actividad
    actividad.verificar_informes()

    return redirect('listar_informe', actividad_id=actividad.id)

@never_cache    
@login_required
def agregar_informe(request, actividad_id):
    user = request.user
    permissions = get_user_permissions(user)

    # Asegúrate de que la actividad existe y el usuario es colaborador
    actividad = get_object_or_404(Actividad, id=actividad_id)
    
    # Verificar si la solicitud es AJAX para devolver las fechas
    if request.headers.get('X-Requested-With') == 'XMLHttpRequest' and request.method == 'GET':
        min_fecha_inicio = actividad.fecha_inicio.isoformat()
        max_fecha_fin = actividad.fecha_fin.isoformat()
        return JsonResponse({'min_fecha_inicio': min_fecha_inicio, 'max_fecha_fin': max_fecha_fin})

    if request.method == 'POST':
        form = InformeForm(request.POST, request.FILES, actividad=actividad)
        if form.is_valid():
            informe = form.save(commit=False)
            informe.actividad = actividad
            informe.usuario = user  # Asigna el usuario que crea el informe
            informe.save()

            # Crear un registro de firma para cada colaborador
            for colaborador in actividad.colaborador.all():
                print(f"Colaborador: {colaborador.first_name}")
                InformeFirmado.objects.create(informe=informe, usuario=colaborador, estado_firma=False)

            actividad.verificar_informes()
            messages.success(request, 'Informe y firmas creados correctamente.')

            return redirect('listar_informe', actividad_id=actividad.id)
        else:
            messages.error(request, 'Hubo un error con el formulario de informe.')

    else:
        form = InformeForm(actividad=actividad)
        
    actividad.verificar_informes()
    print(f"Fecha inicio de la actividad: {actividad.fecha_inicio.isoformat()}")
    print(f"Fecha fin de la actividad: {actividad.fecha_fin.isoformat()}")

    context = {
        'form': form,
        'actividad': actividad,
        **permissions,
    }

    return render(request, 'actividades/listar_informe.html', context)

def editar_informe_modal(request, informe_id):
    try:
        if request.headers.get('X-Requested-With') == 'XMLHttpRequest' and request.method == 'POST':
            informe = get_object_or_404(Informe, id=informe_id)
    
            titulo = request.POST.get('titulo')
            fecha_inicio = request.POST.get('fecha_inicio')
            fecha_fin = request.POST.get('fecha_fin')
    
            # Validar y actualizar los datos
            if titulo and fecha_inicio and fecha_fin:
                informe.titulo = titulo
                informe.fecha_inicio = fecha_inicio
                informe.fecha_fin = fecha_fin
                informe.save()
                return JsonResponse({'success': True, 'message': 'Informe actualizado correctamente.'})
            else:
                return JsonResponse({'success': False, 'message': 'Faltan datos para actualizar.'}, status=400)
    
        elif request.method == 'GET':  # Enviar datos del informe para el modal
            informe = get_object_or_404(Informe, id=informe_id)
            actividad = informe.actividad  # Asume que existe una relación con la actividad
    
            return JsonResponse({
                'titulo': informe.titulo,
                'fecha_inicio_informe': informe.fecha_inicio,
                'fecha_fin_informe': informe.fecha_fin,
                'fecha_inicio_actividad': actividad.fecha_inicio,
                'fecha_fin_actividad': actividad.fecha_fin,
            })
    
        return JsonResponse({'success': False, 'message': 'Método no permitido.'}, status=405)
    except Exception as e:
            return JsonResponse({'success': False, 'message': f'Error: {str(e)}'}, status=500)

@never_cache
@login_required
def editar_informe(request, informe_id):
    user = request.user
    permissions = get_user_permissions(user)
    informe = get_object_or_404(Informe, id=informe_id)
    firma = get_object_or_404(InformeFirmado, informe=informe, usuario=user)
    perfil_usuario = get_object_or_404(PerfilUsuario, user=user)
    actividad = informe.actividad

    if request.method == 'POST':
        print(request.POST)
        # Inicializa ambos formularios
        form_informe = InformeContenidoForm(request.POST, request.FILES, instance=informe)
        form_evidencia = EvidenciaForm(request.POST, request.FILES)

        # Verifica qué acción se realizó
        if 'action' in request.POST:
            action = request.POST['action']

            if action == 'guardar':
                # Valida y guarda el formulario del informe
                if form_informe.is_valid():
                    form_informe.save()
                    messages.success(request, 'Informe actualizado exitosamente.')
                    
                    # Ahora maneja la evidencia
                    if form_evidencia.is_valid():  # Verifica si el formulario de evidencia es válido
                        # Verifica si se ha subido un archivo en la evidencia
                        if form_evidencia.cleaned_data['imagen'] or form_evidencia.cleaned_data['archivo'] or form_evidencia.cleaned_data['enlace']:
                            nueva_evidencia = form_evidencia.save(commit=False)
                            nueva_evidencia.informe = informe  # Asocia la evidencia al informe
                            nueva_evidencia.save()
                            messages.success(request, 'Evidencia agregada exitosamente.')
                        else:
                            messages.info(request, 'No se agregó ninguna nueva evidencia. Se conservarán las evidencias existentes.')
                    else:
                        messages.error(request, 'Hubo un error al agregar la evidencia.')

                    return redirect('listar_informe', actividad_id=informe.actividad.id)
                else:
                    messages.error(request, 'Hubo un error al actualizar el informe.')

            elif action == 'firmar':
                # Lógica para firmar el informe
                firma, created = InformeFirmado.objects.get_or_create(informe=informe, usuario=request.user)
                if not firma.estado_firma:  # Si aún no ha firmado
                    firma.estado_firma = True
                    firma.save()
                    actividad.verificar_informes()
                    if actividad.estado == 'Finalizada':
                        messages.success(request, 'Todos los informes han sido firmados. Actividad finalizada.')
                    else:
                        messages.success(request, 'Informe firmado correctamente.')
                else:
                    messages.info(request, 'Ya habías firmado este informe.')

    else:
        form_informe = InformeContenidoForm(instance=informe)
        form_evidencia = EvidenciaForm()

    # Obtener evidencias existentes y clasificar los archivos por tipo
    evidencias = informe.evidencias.all()  # Obtener todas las evidencias

    # Clasificar las evidencias según el tipo de archivo
    for evidencia in evidencias:
        if evidencia.archivo:
            if evidencia.archivo.name.endswith('.pdf'):
                evidencia.tipo = 'pdf'
            elif evidencia.archivo.name.endswith('.xlsx'):
                evidencia.tipo = 'excel'
            elif evidencia.archivo.name.endswith('.docx'):
                evidencia.tipo = 'word'
            else:
                evidencia.tipo = 'otros'
        elif evidencia.enlace:
            evidencia.tipo = 'enlace'
        else:
            evidencia.tipo = 'imagen'
            
    actividad.verificar_informes()

    context = {
        'form_informe': form_informe,
        'form_evidencia': form_evidencia,
        'informe': informe,
        'firma_imagen': perfil_usuario.img_firma,
        'evidencias': evidencias,  # Añadimos las evidencias al contexto
        **permissions
    }

    return render(request, 'actividades/editar_informe.html', context)
@never_cache
@login_required
def eliminar_informe(request, informe_id):
    informe = get_object_or_404(Informe, id=informe_id)
    actividad = informe.actividad
    # Verifica si el usuario actual es el creador del informe
    if informe.usuario != request.user:
        messages.error(request, "No tienes permiso para eliminar este informe.")
        return redirect('listar_informe',actividad_id=actividad.id)  # Redirige a la lista de informes o a donde desees

    # Si el usuario es el creador, procede a eliminar el informe
    informe.delete()
    messages.success(request, "Informe eliminado exitosamente.")
    actividad.verificar_informes()
    return redirect('listar_informe',actividad_id=actividad.id)  # Redirige a la lista de informes o a donde desees
@never_cache
@login_required
def eliminar_evidencia(request, evidencia_id):
    evidencia = get_object_or_404(Evidencia, id=evidencia_id)
    # Si deseas eliminar la imagen del sistema de archivos
    if evidencia.imagen:
        evidencia.imagen.delete(save=False)  # Eliminar archivo de imagen

    actividad_id = evidencia.informe.actividad.id  # Obtener ID de la actividad
    evidencia.delete()  # Eliminar el objeto de la base de datos

    return redirect('editar_informe', informe_id=evidencia.informe.id)
@never_cache
@login_required
def firmar_informe(request, informe_id):
    # Obtener el informe a partir del ID
    informe = get_object_or_404(Informe, id=informe_id)

    # Verifica si el usuario ya ha firmado el informe
    firma_existente = InformeFirmado.objects.filter(informe=informe, usuario=request.user)

    if firma_existente.exists():
        messages.warning(request, 'Ya has firmado este informe.')
    else:
        try:
            # Crear un nuevo registro de firma
            firma = InformeFirmado.objects.create(informe=informe, usuario=request.user, estado_firma=True)
            print(f'Registro de firma creado correctamente para {request.user.username}')
            
            # Mensaje de éxito
            messages.success(request, 'Informe firmado exitosamente.')
        except Exception as e:
            print(f'Error al crear registro de firma para {request.user.username}: {e}')
            messages.error(request, 'Hubo un error al firmar el informe.')

    # Obtener la actividad relacionada con el informe
    actividad = informe.actividad  # Asegúrate de que Informe tiene un campo de relación con Actividad

    return redirect('listar_informe', actividad_id=actividad.id)

@never_cache#############################################                                TIPO DE ACTIVIDAD                      #################################################################
@login_required
def listar_tipos(request):
    user = request.user  
    permissions = get_user_permissions(user)
    
    perfil_usuario = PerfilUsuario.objects.get(user=request.user)
    programa_usuario = perfil_usuario.programa

    # Verificar si es una petición AJAX
    if request.headers.get('x-requested-with') == 'XMLHttpRequest':
        query = request.GET.get('q', '').strip()
        tipos = (
            TipoActividad.objects.filter(nombre__icontains=query, programa=programa_usuario)
            if query else
            TipoActividad.objects.filter(programa=programa_usuario)
        )
        tipos_data = [{'id': tipo.id,'nombre': tipo.nombre} for tipo in tipos]
        return JsonResponse({'tipos': tipos_data})


    # Si no es AJAX, manejar la lógica normal
    tipos = TipoActividad.objects.filter(programa=programa_usuario)
    paginator = Paginator(tipos, 10)  # 10 tipos por página
    page_number = request.GET.get('page')
    page_obj = paginator.get_page(page_number)

    context = {
        'tipos': page_obj.object_list,
        'page_obj': page_obj,
        **permissions
    }
    return render(request, 'actividades/listar_tipoActividad.html', context)
@never_cache
@login_required
def agregar_tipo(request):
    if request.method == 'POST':
        form = TipoActividadForm(request.POST)
        print(f"Datos recibidos en el formulario: {request.POST}")
        
        # Obtener el perfil del usuario autenticado
        perfil_usuario = PerfilUsuario.objects.get(user=request.user)
        print(f"Perfil del usuario autenticado: {perfil_usuario}")
        
        # Asignar el programa del usuario al tipo de actividad
        if perfil_usuario.programa:
            print(f"Programa del perfil del usuario: {perfil_usuario.programa}")
            # Si el formulario es válido, guardamos el objeto y asignamos el programa
            if form.is_valid():
                tipo_actividad = form.save(commit=False)  # No guardamos aún el objeto
                tipo_actividad.programa = perfil_usuario.programa  # Asignamos el programa del usuario
                tipo_actividad.save()  # Ahora guardamos el objeto con el programa asignado
                
                print("Tipo de actividad guardado correctamente")
                return JsonResponse({'success': True})
            else:
                print("Formulario no válido")
                return JsonResponse({'success': False, 'message': 'Formulario no válido'})
        else:
            return JsonResponse({'success': False, 'message': 'El usuario no tiene un programa asignado'})
        
    return JsonResponse({'success': False, 'message': 'Método no permitido'})
@never_cache
@login_required
def editar_tipo(request):
    if request.method == 'POST':
        tipo_id = request.POST.get('tipo_id')
        tipo = get_object_or_404(TipoActividad, id=tipo_id)
        form = TipoActividadForm(request.POST, instance=tipo)
        if form.is_valid():
            form.save()
            return JsonResponse({'success': True})
    return JsonResponse({'success': False})
@never_cache
@login_required
def eliminar_tipo(request, tipo_id):
    tipo = get_object_or_404(TipoActividad, id=tipo_id)
    tipo.delete()
    return redirect('listar_tipos')


@never_cache#############################################                                   USUARIO                             #########################################################
@login_required
def listar_usuarios(request):
    user = request.user  
    permissions = get_user_permissions(user)
    
    if request.headers.get('x-requested-with') == 'XMLHttpRequest':
        query = request.GET.get('q', '').strip()
        
        # Buscar usuarios por `first_name` o `last_name`
        usuarios = User.objects.filter(
            perfilusuario__programa=user.perfilusuario.programa
        ).filter(
            Q(first_name__icontains=query) | Q(last_name__icontains=query)
        ).exclude(Q(username='admin') | Q(groups__name='Director')) if query else User.objects.filter(
            perfilusuario__programa=user.perfilusuario.programa
        ).exclude(Q(username='admin') | Q(groups__name='Director'))
        
        # Serializar los datos de los usuarios
        usuarios_data = [{
            'id': usuario.id,
            'username': usuario.username,
            'password': usuario.password,
            'first_name': usuario.first_name,
            'last_name': usuario.last_name,
            'email': usuario.email,
            'groups': [{'id': group.id, 'name': group.name} for group in usuario.groups.all()]
        } for usuario in usuarios]
        
        return JsonResponse({'usuarios': usuarios_data})
    
    # Lógica normal para solicitudes no AJAX
    perfil_usuario = user.perfilusuario
    usuarios = User.objects.filter(perfilusuario__programa=perfil_usuario.programa).exclude(Q(username='admin') | Q(groups__name='Director'))
    paginator = Paginator(usuarios, 15)
    page_number = request.GET.get('page')
    page_obj = paginator.get_page(page_number)
    
    form_usuario = UserManagementForm()
    context = {
        'page_obj': page_obj,
        'usuarios': page_obj.object_list,
        'form_usuario': form_usuario,
        **permissions
    }
    return render(request, 'actividades/listar_usuario.html', context)
@never_cache
@login_required
def agregar_usuario(request):
    user = request.user
    if request.method == "POST":
        form = UserManagementForm(request.POST)
        if form.is_valid():
            user = form.save(commit=False)  # No guardar aún
            user.set_password(form.cleaned_data['password'])  # Asegurarse de encriptar la contraseña
            user.save()  # Ahora guardar el usuario
            group = form.cleaned_data['group']  # Obtener el grupo seleccionado
            group.user_set.add(user)  # Asignar el usuario al grupo
            
            # Obtener el perfil del creador
            if request.user.perfilusuario:  # Asegurarse de que el creador tiene un perfil
                creator_profile = request.user.perfilusuario
                # Aquí asignamos la facultad y el programa del creador al nuevo usuario
                user.perfilusuario.programa = creator_profile.programa
                user.perfilusuario.save()

            messages.success(request, "Usuario agregado exitosamente.")
            return redirect('listar_usuarios')  # Redirigir después de agregar el usuario
        else:
            messages.error(request, "Error al agregar usuario. Revisa los datos.")
    return redirect('listar_usuarios')  # Redirigir si el formulario no es válido
@never_cache
@login_required
def cargar_usuarios_csv(request):
    if request.method == 'POST':
        # Verificar si se subió un archivo CSV
        archivo_csv = request.FILES.get('archivo_csv')

        if not archivo_csv:
            messages.error(request, "Por favor, selecciona un archivo CSV.")
            return render(request, 'cargar_usuarios_csv.html')

        try:
            # Leer el archivo CSV usando pandas
            df = pd.read_csv(archivo_csv)

            # Verificar que las columnas necesarias estén presentes
            columnas_requeridas = ['Cedula', 'Nombres', 'Apellidos', 'Correo', 'Rol']
            if not all(col in df.columns for col in columnas_requeridas):
                messages.error(request, "El archivo CSV debe contener las columnas: 'Cedula','Nombres','Apellidos','Correo', 'Rol'.")
                return render(request, 'listar_usuarios.html')
            
            # Obtener el perfil del usuario que está realizando la carga
            creator_profile = getattr(request.user, 'perfilusuario', None)

            # Crear usuarios a partir del CSV
            for _, row in df.iterrows():
                # Crear el usuario solo si no existe un usuario con el mismo nombre de usuario
                if not User.objects.filter(username=row['Cedula']).exists():
                    user = User.objects.create_user(
                        username=row['Cedula'],
                        first_name=row['Nombres'],
                        last_name=row['Apellidos'],
                        email=row['Correo'],
                        password=row['Correo']
                    )
                    # Asignar el usuario al grupo especificado
                    grupo, created = Group.objects.get_or_create(name=row['Rol'])
                    user.groups.add(grupo)
                    
                    if request.user.perfilusuario:  # Asegurarse de que el creador tiene un perfil
                        creator_profile = request.user.perfilusuario
                        # Aquí asignamos la facultad y el programa del creador al nuevo usuario
                        user.perfilusuario.programa = creator_profile.programa
                        user.perfilusuario.save()
                    

            messages.success(request, "Usuarios creados exitosamente.")
            return redirect('listar_usuarios')
        except Exception as e:
            messages.error(request, f"Error al procesar el archivo CSV: {str(e)}")

    return render(request, 'actividades/listar_usuario.html')
@never_cache
@login_required
def editar_usuario(request):
    if request.method == "POST":
        usuario_id = request.POST.get('usuario_id')
        usuario = get_object_or_404(User, id=usuario_id)
        form = UserManagementForm(request.POST, instance=usuario)
        if form.is_valid():
            user = form.save(commit=False)  
            user.set_password(form.cleaned_data['password']) 
            user.save()  
            group = form.cleaned_data['group']  
            # Limpiar las asignaciones de grupos anteriores
            usuario.groups.clear()  
            group.user_set.add(user)  
            messages.success(request, "Usuario editado exitosamente.")
            return redirect('listar_usuarios')
        else:
            messages.error(request, "Error al editar usuario. Revisa los datos.")
    return redirect('listar_usuarios')
@never_cache
@login_required
def eliminar_usuario(request, usuario_id):
    usuario = get_object_or_404(User, id=usuario_id)
    usuario.delete()
    messages.success(request, "Usuario eliminado exitosamente.")
    return redirect('listar_usuarios')


@never_cache#############################################                                   PERFIL                             #########################################################
@login_required
def perfil(request):
    user = request.user
    permissions = get_user_permissions(user)
    perfil_usuario = PerfilUsuario.objects.get(user=user)

    if request.method == "POST":
        user_form = UserPerfilForm(request.POST, instance=user)
        perfil_form = PerfilUsuarioForm(request.POST, request.FILES, instance=perfil_usuario)

        if user_form.is_valid() and perfil_form.is_valid():
            # Verifica si hay una nueva foto de perfil
            if 'img_usuario' in request.FILES:
                foto_usuario = request.FILES['img_usuario']
                img = PILImage.open(foto_usuario)
                img = img.resize((300, 300), PILImage.LANCZOS)  # Redimensiona la imagen de perfil a 300x300

                if img.mode in ('RGBA', 'P'):
                    img = img.convert('RGB')

                # Guarda la imagen redimensionada en un objeto BytesIO
                thumb_io = BytesIO()
                img.save(thumb_io, format='JPEG')

                # Crea un archivo de imagen a partir del objeto BytesIO
                nuevo_foto_usuario = ContentFile(thumb_io.getvalue(), name=foto_usuario.name)
                perfil_form.instance.img_usuario = nuevo_foto_usuario

            # Verifica si hay una nueva firma
            if 'img_firma' in request.FILES:
                foto_firma = request.FILES['img_firma']
                img_firma = PILImage.open(foto_firma)
                img_firma = img_firma.resize((600, 200), PILImage.LANCZOS)  # Redimensiona la firma a 600x200

                if img_firma.mode in ('RGBA', 'P'):
                    img_firma = img_firma.convert('RGB')

                # Guarda la imagen de la firma redimensionada
                thumb_io_firma = BytesIO()
                img_firma.save(thumb_io_firma, format='JPEG')

                nuevo_foto_firma = ContentFile(thumb_io_firma.getvalue(), name=foto_firma.name)
                perfil_form.instance.img_firma = nuevo_foto_firma

            # Guarda los formularios
            user_form.save()
            perfil_form.save()

            return redirect('dashboard')  # Redirige después de guardar los cambios
        else:
            # Imprime errores para depuración
            print(user_form.errors)
            print(perfil_form.errors)

    else:
        user_form = UserPerfilForm(instance=user)
        perfil_form = PerfilUsuarioForm(instance=perfil_usuario)

    return render(request, "actividades/perfil.html", {
        'user_form': user_form,
        'perfil_form': perfil_form,
        'user': user,
        **permissions
    })

#############################################                                   REPORTES                             #########################################################
@never_cache
@login_required
def generar_informe_pdf(request):
    user = request.user
    permissions = get_user_permissions(user)

    # Identificar tipo de usuario
    es_director = user.groups.filter(name='director').exists()
    es_profesor = user.groups.filter(name='profesor').exists()

    # Filtros
    fecha_inicio = request.GET.get('fecha_inicio')
    fecha_fin = request.GET.get('fecha_fin')
    estado = request.GET.get('estado')

    # Obtener actividades según el tipo de usuario
    if es_director:
        actividades = Actividad.objects.filter(
            creado_por__perfilusuario__programa=user.perfilusuario.programa
        ).distinct()
    elif es_profesor:
        actividades = Actividad.objects.filter(
            Q(creado_por=user) | Q(colaborador=user)
        ).distinct()
    else:
        actividades = Actividad.objects.none()

    # Filtro por estado
    if estado:
        actividades = actividades.filter(estado__iexact=estado)  # Comparación sin importar mayúsculas/minúsculas

    # Filtro por fechas
    try:
        if fecha_inicio:
            fecha_inicio = datetime.strptime(fecha_inicio, '%Y-%m-%d').date()
            actividades = actividades.filter(fecha_inicio__gte=fecha_inicio)

        if fecha_fin:
            fecha_fin = datetime.strptime(fecha_fin, '%Y-%m-%d').date()
            actividades = actividades.filter(fecha_fin__lte=fecha_fin)
    except ValueError:
        pass  # Si las fechas no son válidas, simplemente no aplicamos el filtro

    # Contexto para el template
    context = {
        'usuario': user,
        'actividades': actividades,
        'fecha_inicio': fecha_inicio,
        'fecha_fin': fecha_fin,
        'estado': estado,
        **permissions
    }

    # Renderizar HTML
    template = get_template('actividades/reporte_actividades.html')
    html = template.render(context)

    # Generar PDF
    response = HttpResponse(content_type='application/pdf')
    response['Content-Disposition'] = 'inline; filename="informe_actividades.pdf"'

    pisa_status = pisa.CreatePDF(html, dest=response)
    if pisa_status.err:
        return HttpResponse('Error al generar el PDF', status=500)

    return response

class DescargarReporte(View):
    def get(self, request, actividad_id, formato):
        actividad = Actividad.objects.get(id=actividad_id)
        informes = Informe.objects.filter(actividad=actividad)

        if formato == 'pdf':
            return self.generar_pdf(actividad, informes)
        elif formato == 'word':
            return self.generar_word(actividad, informes)
        elif formato == 'ppt':
            return self.generar_ppt(actividad, informes)
        else:
            return HttpResponse("Formato no soportado", status=400)

    def generar_pdf(self, actividad, informes):
        response = HttpResponse(content_type='application/pdf')
        response['Content-Disposition'] = f'inline; filename="reporte_actividad_{actividad.id}.pdf"'
        
        # Crear documento PDF
        doc = SimpleDocTemplate(response, pagesize=letter)
        # Establecer los metadatos, incluido el título
        doc.title = f"Reporte de Actividad {actividad.id}"
        styles = getSampleStyleSheet()
        flowables = [] 
        centered_heading_style = ParagraphStyle(
            name="CenteredHeading",
            parent=styles['Heading2'],
            alignment=TA_CENTER
        )
        
        # Definir el tamaño de la página
        PAGE_HEIGHT = letter[1]
        
        # Portada
        logo_path = "actividades/static/actividades/img/cesmag.png"
        flowables.append(Image(logo_path, width=2 * inch, height=2 * inch))
        
        flowables.append(Spacer(1, 0.5 * inch))  # Espaciado
        flowables.append(Paragraph("Informe de Actividad", styles['Title']))
        
        flowables.append(Spacer(1, 0.1 * inch))
        flowables.append(Paragraph(actividad.titulo, styles['Title']))
        
        flowables.append(Spacer(1, 0.5 * inch))
        flowables.append(Paragraph("Nombre del Docente:", centered_heading_style))

        colaboradores = actividad.colaborador.all()
        nombres_colaboradores = ', '.join([f"{colaborador.first_name} {colaborador.last_name}" for colaborador in colaboradores])
        flowables.append(Paragraph(nombres_colaboradores, centered_heading_style))  # Centrando nombres

        flowables.append(Spacer(1, 0.5 * inch))
        flowables.append(Paragraph("Programa Académico:", centered_heading_style))
        programa_academico = "No especificado"  # Valor predeterminado
        if actividad.creado_por:
            perfil_usuario = PerfilUsuario.objects.filter(user=actividad.creado_por).first()
            if perfil_usuario and perfil_usuario.programa:
                programa_academico = perfil_usuario.programa.nombre  # Asumiendo que el programa tiene un campo 'nombre'

        flowables.append(Paragraph(programa_academico, centered_heading_style))

        # Añadir un espaciador dinámico para empujar la fecha al final
        flowables.append(Spacer(1, PAGE_HEIGHT - len(flowables) * 60))  # Ajusta este número según el contenido

        # Añadir la fecha en la parte inferior
        flowables.append(Paragraph(f"Fecha: {datetime.now().strftime('%Y-%m-%d')}", centered_heading_style))
        
        flowables.append(PageBreak())  # Nueva página para los informes

        PAGE_HEIGHT = letter[1]
        
        # Estilos para los párrafos
        styles = getSampleStyleSheet()
        paragraph_style = styles['Normal']
        
        # Sección de Informes
        for i, informe in enumerate(informes):
            # Crear la tabla con los datos del informe
            data = [
                [Paragraph(informe.titulo, styles['Normal'])],  # Título del informe
                ['Fecha inicio:',Paragraph(informe.fecha_inicio.strftime("%d/%m/%Y"), styles['Normal']), 'Fecha fin:',Paragraph(informe.fecha_fin.strftime("%d/%m/%Y")) ],  # Encabezado de fechas
                [Paragraph('Descripción', styles['Normal'])],  # Descripción
                [Paragraph(informe.contenido, paragraph_style)],  # Contenido largo
                ['Evidencias'],  # Título para la sección de evidencias
            ]

            # Crear una lista de imágenes en una sola fila
            imagenes = []
            imagen_fila = []
            documentos_enlaces = []
            hay_evidencias = False
            titulo_agregado = False
            max_ancho_tabla = 7 * inch  # Ancho total de la tabla (4 columnas de 2 pulgadas)

            # Ancho máximo de cada imagen para que encajen bien
            ancho_imagen = max_ancho_tabla / 3  # 3 imágenes por fila

            # Agregar las imágenes a la fila
            for evidencia in informe.evidencias.all():
                if evidencia.imagen:  # Verifica si 'imagen' tiene un archivo asociado
                    imagen = Image(evidencia.imagen.path, width=ancho_imagen, height=ancho_imagen * 0.75)  # Ajustar la proporción
                    imagen_fila.append(imagen)

                    # Si ya tenemos 3 imágenes en la fila, la agregamos a la lista de filas y empezamos una nueva fila
                    if len(imagen_fila) == 3:
                        imagenes.append(imagen_fila)
                        imagen_fila = []  # Reiniciar la fila

            # Si quedaron imágenes en la última fila, agrégalas
            if imagen_fila:
                imagenes.append(imagen_fila)
            # Crear una tabla para las imágenes solo si hay imágenes
            if imagenes:
            # Crear una tabla para las imágenes, dentro de la celda de evidencias
                tabla_imagenes = Table(imagenes, colWidths=[ancho_imagen] * 3)  # 3 columnas por fila
                tabla_imagenes.setStyle(TableStyle([
                    ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),  # Centrar verticalmente las imágenes
                    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),  # Centrar horizontalmente las imágenes
                ]))
                # Agregar la tabla de imágenes a los datos
                data.append([tabla_imagenes])
            
            # Agregar título y documentos/enlaces en una subtabla
            titulo_agregado = False  # Control para agregar el título solo una vez
            for evidencia in informe.evidencias.all():
                if evidencia.archivo or evidencia.enlace:
                    hay_evidencias = True

                    # Agregar título de "Documentos:" solo una vez
                    if not titulo_agregado:
                        documentos_enlaces.append([Paragraph("Documentos:", styles['Normal'])])
                        titulo_agregado = True

                    # Agregar cada documento o enlace en la subtabla
                    if evidencia.archivo:
                        documentos_enlaces.append([Paragraph(f"<a href='{evidencia.archivo.url}' color='green'>{evidencia.archivo}</a>", styles['Normal'])])
                    elif evidencia.enlace:
                        documentos_enlaces.append([Paragraph(f"<a href='{evidencia.enlace}' color='blue'>{evidencia.enlace}</a>", styles['Normal'])])

            # Si hay documentos o enlaces, los agregamos en una subtabla dentro de la celda principal
            if documentos_enlaces:
                tabla_documentos_enlaces = Table(documentos_enlaces)
                tabla_documentos_enlaces.setStyle(TableStyle([
                    ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                    ('ALIGN', (0, 0), (-1, -1), 'LEFT'),  # Alineado a la izquierda para los enlaces y documentos
                ]))
                data.append([tabla_documentos_enlaces])

            # Si no hay ningún tipo de evidencia, mostrar el mensaje de "No hay evidencias disponibles"
            if not hay_evidencias:
                data.append([Paragraph("No hay evidencias disponibles.", styles['Normal'])])
        
            # Crear la tabla
            table = Table(data, colWidths=[1.8 * inch, 1.8 * inch, 1.8 * inch, 1.8 * inch])

            # Establecer el estilo de la tabla
            table.setStyle(TableStyle([
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.black),
                ('ALIGN', (0, 0), (-1, -1), 'LEFT'),  # Alinear texto a la izquierda
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('GRID', (0, 0), (-1, -1), 1, colors.black),
                ('VALIGN', (0, 0), (-1, -1), 'TOP'),  # Alinear texto al inicio de la celda
            ]))

            # Ajustar el colspan para el título y descripción
            table.setStyle(TableStyle([
                ('SPAN', (0, 0), (3, 0)),  # Combinar celdas para el título
                ('SPAN', (0, 2), (3, 2)),  # Combinar celdas para la descripción
                ('SPAN', (0, 3), (3, 3)),  # Combinar celdas para el contenido largo
                ('SPAN', (0, 4), (3, 4)),  # Combinar celdas para el título "Evidencias"
                ('SPAN', (0, 5), (3, 5)),  # Combinar celdas para el título "Evidencias"
                ('SPAN', (0, 6), (3, 6)),  # Combinar celdas para el título "Archivos"
                ('SPAN', (0, 7), (3, 7)),  # Combinar celdas para el título "Enlaces"
                
            ]))

            # Agregar la tabla al flujo
            flowables.append(table)
            
            # Calcular el espacio ocupado por la tabla de evidencias
            espacio_ocupado = table.wrap(doc.width, doc.height)[1]  # Obtener altura ocupada por la tabla

            # Calcular el espacio restante en la página
            espacio_restante = PAGE_HEIGHT - (len(flowables) * 0.5 * inch) - espacio_ocupado
            # Agregar un salto de página solo si no es la última iteración
            if i < len(informes) - 1:
                flowables.append(Spacer(1, 24))
            
        # Obtener los estilos
        styles = getSampleStyleSheet()
        centered_heading_style = styles['Heading3']

        # Inicializar la lista de filas para la tabla de firmas
        filas_firmas = []
        fila_actual = []

        # Número máximo de columnas en la tabla
        max_columnas = 3

        for colaborador in colaboradores:
            try:
                perfil_usuario = colaborador.perfilusuario

                # Crear celda con firma o mensaje "Sin firma"
                if perfil_usuario.img_firma and perfil_usuario.img_firma.path:
                    nombre_apellido = Paragraph(f"{colaborador.first_name} {colaborador.last_name}", centered_heading_style)
                    firma_imagen = Image(perfil_usuario.img_firma.path, width=1.5 * inch, height=0.75 * inch)
                    celda = [firma_imagen, nombre_apellido]
                else:
                    nombre_apellido = Paragraph(f"{colaborador.first_name} {colaborador.last_name}", centered_heading_style)
                    firma_imagen = Paragraph("Sin firma", centered_heading_style)
                    celda = [firma_imagen, nombre_apellido]

                # Agregar la celda a la fila actual
                fila_actual.append(celda)

                # Añadir fila completa a la tabla si llega al máximo de columnas
                if len(fila_actual) == max_columnas:
                    filas_firmas.append(fila_actual)
                    fila_actual = []

            except PerfilUsuario.DoesNotExist:
                nombre_apellido = Paragraph(f"{colaborador.first_name} {colaborador.last_name}", centered_heading_style)
                firma_imagen = Paragraph("Sin firma", centered_heading_style)
                celda = [firma_imagen, nombre_apellido]
                fila_actual.append(celda)

                if len(fila_actual) == max_columnas:
                    filas_firmas.append(fila_actual)
                    fila_actual = []

        # Agregar la última fila, si quedó incompleta
        if fila_actual:
            filas_firmas.append(fila_actual)

        # Crear la tabla para las firmas solo si hay filas
        if filas_firmas:
            # Crear la tabla para las firmas
            tabla_firmas = Table(filas_firmas, colWidths=[2 * inch] * max_columnas)
            tabla_firmas.setStyle(TableStyle([
                ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ]))

            # Calcular el espacio utilizado por la última tabla (informes)
            espacio_informe = table.wrap(doc.width, doc.height)[1]
            espacio_firmas = tabla_firmas.wrap(doc.width, doc.height)[1]
            espacio_restante = doc.height - espacio_informe # Calcular el espacio restante en la página
            
            altura_por_fila = espacio_firmas / len(filas_firmas) if len(filas_firmas) > 0 else 0

            # Comprobamos si hay suficiente espacio para la tabla completa
            if espacio_restante >= espacio_firmas:
                # Si hay suficiente espacio para la tabla completa, la agregamos sin cambios
                flowables.append(Spacer(1, 12))
                flowables.append(tabla_firmas)

            # Si no hay suficiente espacio para la tabla completa, verificamos si se puede ajustar
            elif espacio_restante >= altura_por_fila:
                # Ajustamos la altura de las filas para que quepan en el espacio restante
                filas_que_caben = int(espacio_restante // altura_por_fila)

                # Si hay filas que caben, ajustamos la tabla
                tabla_firmas_ajustada = Table(filas_firmas[:filas_que_caben], colWidths=[2 * inch] * max_columnas)
                tabla_firmas_ajustada.setStyle(TableStyle([
                    ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                    ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                ]))

                # Añadir la tabla ajustada
                flowables.append(tabla_firmas_ajustada)

                # Si hay más filas, hacer un salto de página y agregar las filas restantes
                if filas_que_caben < len(filas_firmas):
                    flowables.append(PageBreak())  # Salto de página
                    tabla_firmas_restante = Table(filas_firmas[filas_que_caben:], colWidths=[2 * inch] * max_columnas)
                    tabla_firmas_restante.setStyle(TableStyle([
                        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                    ]))
                    flowables.append(tabla_firmas_restante)

            # Si no cabe ni una fila de la tabla, saltamos a la siguiente página
            else:
                # Salto de página y agregamos la tabla completa
                flowables.append(PageBreak())
                flowables.append(tabla_firmas)
        else:
            flowables.append(Paragraph("No hay colaboradores disponibles.", styles['Normal']))

        # Construir el PDF
        doc.build(flowables)
        return response
    
    def generar_word(self, actividad, informes):
        # Crear el documento Word
        doc = Document()
        
        # Establecer márgenes de 2 cm
        sections = doc.sections
        for section in sections:
            section.left_margin = Inches(0.79)  # 2 cm
            section.right_margin = Inches(0.79)  # 2 cm
            section.top_margin = Inches(0.79)
            section.bottom_margin = Inches(0.79)
            
        # Crear una sección para la portada
        section = doc.sections[0]
        section.start_type = 0  # Comenzar en la misma página
        section.different_first_page_header_footer = True

        # Añadir el logo de forma centrada
        logo_path = "actividades/static/actividades/img/cesmag.png"
        doc.add_picture(logo_path, width=Inches(2), height=Inches(2))
        doc.paragraphs[-1].alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        
        doc.add_paragraph("\n" * 1)

        # Información del informe (centrada)
        doc.add_paragraph("Informe de Actividad", style='Heading1').alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        doc.add_paragraph(actividad.titulo, style='Heading2').alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        
        doc.add_paragraph("\n" * 3)
        
        doc.add_paragraph("Nombre del Profesor:" ,style='Heading2').alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        colaboradores = actividad.colaborador.all()
        nombres_colaboradores = ', '.join([f"{colaborador.first_name} {colaborador.last_name}" for colaborador in colaboradores])
        doc.add_paragraph(nombres_colaboradores).alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        
        # Añadir el programa académico
        programa_academico = "No especificado"  # Valor predeterminado
        if actividad.creado_por:
            perfil_usuario = PerfilUsuario.objects.filter(user=actividad.creado_por).first()
            if perfil_usuario and perfil_usuario.programa:
                programa_academico = perfil_usuario.programa.nombre  # Ajusta si el campo no es 'nombre'
        doc.add_paragraph("\n" * 1)

        doc.add_paragraph("Programa Académico:", style='Heading2').alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        doc.add_paragraph(programa_academico).alignment = WD_PARAGRAPH_ALIGNMENT.CENTER

        doc.add_paragraph("\n" * 10)
        
        # Fecha de generación del reporte (en la parte inferior de la hoja)
        paragraph = doc.add_paragraph(f"Fecha: {datetime.now().strftime('%Y-%m-%d')}", style='Heading2')
        paragraph.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        
        # Forzar la fecha al final con un espaciado específico
        paragraph_format = paragraph.paragraph_format
        paragraph_format.space_before = Pt(0)
        paragraph_format.space_after = Pt(0)
        paragraph_format.line_spacing = Pt(12)

        # Crear una nueva sección para los informes
        doc.add_page_break()
        doc.add_heading("Informes", level=2)

        # Recorrer y agregar los informes en tablas dinámicas
        for informe in informes:
            # Crear la tabla para cada informe
            table = doc.add_table(rows=0, cols=4)
            table.style = 'Table Grid'

            # Título del informe (Primera fila de la tabla)
            row_cells = table.add_row().cells
            merged_cell = row_cells[0].merge(row_cells[3])
            merged_cell.text = informe.titulo

            # Fechas (Segunda fila de la tabla)
            row_cells = table.add_row().cells
            row_cells[0].text = "Fecha de Inicio"
            row_cells[1].text = informe.fecha_inicio.strftime('%d/%m/%Y')
            row_cells[2].text = "Fecha de Fin"
            row_cells[3].text = informe.fecha_fin.strftime('%d/%m/%Y')
            

            row_cells = table.add_row().cells
            merged_cell = row_cells[0].merge(row_cells[3])
            merged_cell.text = 'Descripcion'
            
            row_cells = table.add_row().cells
            merged_cell = row_cells[0].merge(row_cells[3])
            # Limpiar las etiquetas HTML del contenido
            clean_content = BeautifulSoup(informe.contenido, "html.parser").get_text()

            # Asignar el contenido limpio a la celda
            merged_cell.text = clean_content

            # Evidencias (Cuarta fila de la tabla)
            row_cells = table.add_row().cells
            merged_cell = row_cells[0].merge(row_cells[3])
            merged_cell.text = "Evidencias"
            
            # Contenido de las evidencias (Sexta fila de la tabla)
            row_cells = table.add_row().cells
            merged_cell = row_cells[0].merge(row_cells[3])  # Combinar toda la fila

            images_table = None
            # Construir el contenido de evidencias
            if informe.evidencias.exists():
                # Crear una tabla para las imágenes (2 columnas)
                images_table = merged_cell.add_table(rows=0, cols=2)
                if images_table:
                    tbl = images_table._element
                    tbl.tblPr.append(parse_xml(r'<w:tblBorders %s><w:top w:val="none"/><w:left w:val="none"/><w:bottom w:val="none"/><w:right w:val="none"/><w:insideH w:val="none"/><w:insideV w:val="none"/></w:tblBorders>' % nsdecls('w')))

                # Añadir imágenes a la tabla
                image_row_cells = None
                for idx, evidencia in enumerate(informe.evidencias.filter(imagen__isnull=False)):
                    # Verificar si la imagen tiene un archivo asociado
                    if evidencia.imagen and os.path.exists(evidencia.imagen.path):
                        if idx % 2 == 0:  # Crear una nueva fila cada 2 imágenes
                            image_row_cells = images_table.add_row().cells
                        paragraph = image_row_cells[idx % 2].paragraphs[0]
                        run = paragraph.add_run()
                        run.add_picture(evidencia.imagen.path, width=Inches(3), height=Inches(1.5))
                        
                # Listar archivos y enlaces
                for evidencia in informe.evidencias.all():
                    if evidencia.archivo:
                        merged_cell.add_paragraph(f"Archivo: http://127.0.0.1:8000{evidencia.archivo.url}")
                    if evidencia.enlace:
                        merged_cell.add_paragraph(f"Enlace: {evidencia.enlace}")
            else:
                # Si no hay evidencias, indicar que no hay disponibles
                merged_cell.add_paragraph("No hay evidencias disponibles.")
            doc.add_paragraph("")

        # Crear una tabla de firmas con una sola fila dividida en tres columnas
        signature_table = doc.add_table(rows=1, cols=3)
        signature_table.autofit = True
        signature_table.allow_autofit = True
        # Ajustar ancho dinámico a los márgenes
        for column in table.columns:
            column.width = Cm(2)
        if images_table:
            tbl = images_table._element
            tbl.tblPr.append(parse_xml(r'<w:tblBorders %s><w:top w:val="none"/><w:left w:val="none"/><w:bottom w:val="none"/><w:right w:val="none"/><w:insideH w:val="none"/><w:insideV w:val="none"/></w:tblBorders>' % nsdecls('w')))
        else:
            # Si no hay imágenes, se puede agregar un mensaje o hacer algo por defecto
            pass       

        # Añadir celdas con firmas de los colaboradores
        for idx, colaborador in enumerate(colaboradores):
            perfil_usuario = getattr(colaborador, 'perfilusuario', None)
            row_cells = signature_table.rows[0].cells

            

            if perfil_usuario and perfil_usuario.img_firma:
                row_cells[idx].paragraphs[0].add_run().add_picture(perfil_usuario.img_firma.path, width=Inches(1.5), height=Inches(0.75))
            else:
                row_cells[idx].paragraphs[0].add_run("Sin firma")
            # Agregar el nombre abajo
            row_cells[idx].add_paragraph(f"{colaborador.first_name} {colaborador.last_name}")

        # Guardar el documento en la respuesta HTTP
        response = HttpResponse(content_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document')
        response['Content-Disposition'] = f'attachment; filename="reporte_actividad_{actividad.id}.docx"'
        doc.save(response)
        return response

    

   
    

    def generar_ppt(self, actividad, informes):
        # Crear presentación PowerPoint
        prs = Presentation()

        # Agregar una diapositiva de título
        slide_layout = prs.slide_layouts[0]  # Layout de título
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1] if len(slide.placeholders) > 1 else None  # Verificar si el placeholder existe

        # Añadir el logo centrado en la diapositiva de título
        logo_path = "actividades/static/actividades/img/cesmag.png"
        slide.shapes.add_picture(logo_path, Inches(4), Inches(1), width=Inches(2), height=Inches(2))

        # Añadir el título y el subtítulo
        title.text = "Informe de Actividad"
        if subtitle:
            subtitle.text = actividad.titulo

        # Añadir la información general (centrada)
        slide_layout = prs.slide_layouts[1]  # Layout de título + contenido
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Información del Informe"

        content = slide.shapes.placeholders[1] if len(slide.placeholders) > 1 else slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(7), Inches(4))
        content.text = f"Nombre del Profesor: {', '.join([f'{colaborador.first_name} {colaborador.last_name}' for colaborador in actividad.colaborador.all()])}\n"

        # Información del programa académico
        programa_academico = "No especificado"
        if actividad.creado_por:
            perfil_usuario = PerfilUsuario.objects.filter(user=actividad.creado_por).first()
            if perfil_usuario and perfil_usuario.programa:
                programa_academico = perfil_usuario.programa.nombre

        content.text += f"Programa Académico: {programa_academico}\n"

        # Fecha de creación del reporte
        content.text += f"Fecha: {datetime.now().strftime('%Y-%m-%d')}"

        # Añadir una diapositiva de informes
        for informe in informes:
            slide_layout = prs.slide_layouts[1]  # Layout de título + contenido
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = informe.titulo

            # Verificar si el placeholder para el texto existe, si no, agregar uno
            textbox = slide.shapes.placeholders[1] if len(slide.placeholders) > 1 else slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(7), Inches(4))
            textbox.text = f"Fecha de Inicio: {informe.fecha_inicio.strftime('%d/%m/%Y')}\n"
            textbox.text += f"Fecha de Fin: {informe.fecha_fin.strftime('%d/%m/%Y')}\n"

            # Descripción del informe
            clean_content = BeautifulSoup(informe.contenido, "html.parser").get_text()
            textbox.text += f"Descripción: {clean_content}\n"

            # Añadir las evidencias
            if informe.evidencias.exists():
                textbox.text += "Evidencias:\n"
                for evidencia in informe.evidencias.all():
                    if evidencia.imagen and os.path.exists(evidencia.imagen.path):
                        slide.shapes.add_picture(evidencia.imagen.path, Inches(1), Inches(1.5), width=Inches(3), height=Inches(1.5))
                    if evidencia.archivo:
                        textbox.text += f"Archivo: {evidencia.archivo.url}\n"
                    if evidencia.enlace:
                        textbox.text += f"Enlace: {evidencia.enlace}\n"
            else:
                textbox.text += "No hay evidencias disponibles.\n"

            # Añadir una separación de sección (línea)
            slide.shapes.add_shape(
                MSO_SHAPE, Inches(0.5), Inches(4), Inches(9), Inches(0)  # Agregar una línea
            )

        # Añadir una diapositiva final con las firmas
        slide_layout = prs.slide_layouts[1]  # Layout de título + contenido
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Firmas de los Colaboradores"

        signature_content = slide.shapes.placeholders[1] if len(slide.placeholders) > 1 else slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(7), Inches(4))
        for colaborador in actividad.colaborador.all():
            perfil_usuario = getattr(colaborador, 'perfilusuario', None)
            if perfil_usuario and perfil_usuario.img_firma:
                signature_content.text += f"{colaborador.first_name} {colaborador.last_name}\n"
                signature_content.text += f"Firma: [Imagen de firma]\n"
            else:
                signature_content.text += f"{colaborador.first_name} {colaborador.last_name}\n"
                signature_content.text += "Firma: Sin firma\n"

        # Guardar la presentación como archivo
        pptx_path = f'Informe_actividad_{actividad.id}.pptx'
        prs.save(pptx_path)

        return pptx_path




